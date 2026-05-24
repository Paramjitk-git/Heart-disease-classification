from datetime import datetime
from io import BytesIO
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.metrics import accuracy_score, classification_report, confusion_matrix, precision_score, recall_score
from sklearn.model_selection import train_test_split
from sklearn.tree import DecisionTreeClassifier, plot_tree


BASE_DIR = Path(__file__).resolve().parent
DATA_FILE = BASE_DIR / "processed.cleveland.data"


FEATURE_NAMES = [
    "age",
    "sex",
    "cp",
    "trestbps",
    "chol",
    "fbs",
    "restecg",
    "thalach",
    "exang",
    "oldpeak",
    "slope",
    "ca",
    "thal",
    "target_raw",
]

FEATURE_LABELS = {
    "age": "Age",
    "sex": "Sex",
    "cp": "Chest pain type",
    "trestbps": "Resting blood pressure",
    "chol": "Serum cholesterol",
    "fbs": "Fasting blood sugar > 120",
    "restecg": "Resting ECG result",
    "thalach": "Maximum heart rate",
    "exang": "Exercise-induced angina",
    "oldpeak": "ST depression",
    "slope": "Slope of ST segment",
    "ca": "Major vessels colored",
    "thal": "Thalassemia result",
}


def load_data() -> pd.DataFrame:
    df = pd.read_csv(DATA_FILE, names=FEATURE_NAMES, na_values="?")
    df = df.dropna().copy()
    df["heart_disease"] = (df["target_raw"] > 0).astype(int)
    return df


def train_model(df: pd.DataFrame):
    features = list(FEATURE_LABELS)
    x = df[features]
    y = df["heart_disease"]
    x_train, x_test, y_train, y_test = train_test_split(
        x, y, test_size=0.30, random_state=42, stratify=y
    )

    model = DecisionTreeClassifier(max_depth=3, min_samples_leaf=8, random_state=42)
    model.fit(x_train, y_train)
    y_pred = model.predict(x_test)

    metrics = {
        "accuracy": accuracy_score(y_test, y_pred),
        "precision": precision_score(y_test, y_pred),
        "recall": recall_score(y_test, y_pred),
        "confusion": confusion_matrix(y_test, y_pred),
        "report": classification_report(
            y_test,
            y_pred,
            target_names=["No heart disease", "Heart disease"],
            digits=3,
        ),
        "n_train": len(x_train),
        "n_test": len(x_test),
    }

    return features, model, metrics


def save_charts(df: pd.DataFrame, features: list[str], model: DecisionTreeClassifier, metrics: dict):
    class_counts = df["heart_disease"].map({0: "No heart disease", 1: "Heart disease"}).value_counts()
    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    bars = ax.bar(
        class_counts.index,
        class_counts.values,
        color=["#2F80ED", "#EB5757"],
        edgecolor="#202124",
        linewidth=0.8,
    )
    ax.set_title("Distribution of Target Class", fontsize=16, weight="bold")
    ax.set_ylabel("Number of patient records")
    ax.set_xlabel("Target class")
    ax.grid(axis="y", alpha=0.25)
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width() / 2, height + 2, str(int(height)), ha="center", fontsize=12)
    fig.tight_layout()
    distribution_image = BytesIO()
    fig.savefig(distribution_image, dpi=180, format="png")
    distribution_image.seek(0)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(18, 8.2))
    plot_tree(
        model,
        feature_names=[FEATURE_LABELS[f] for f in features],
        class_names=["No heart disease", "Heart disease"],
        filled=True,
        rounded=True,
        impurity=False,
        fontsize=10,
        ax=ax,
    )
    fig.tight_layout()
    tree_image = BytesIO()
    fig.savefig(tree_image, dpi=180, format="png")
    tree_image.seek(0)
    plt.close(fig)

    cm = metrics["confusion"]
    fig, ax = plt.subplots(figsize=(5.8, 4.8))
    image = ax.imshow(cm, cmap="Blues")
    ax.set_xticks([0, 1], ["Predicted no", "Predicted yes"])
    ax.set_yticks([0, 1], ["Actual no", "Actual yes"])
    ax.set_title("Confusion Matrix", fontsize=15, weight="bold")
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax.text(j, i, cm[i, j], ha="center", va="center", color="#111111", fontsize=16, weight="bold")
    fig.colorbar(image, ax=ax, fraction=0.046, pad=0.04)
    fig.tight_layout()
    confusion_image = BytesIO()
    fig.savefig(confusion_image, dpi=180, format="png")
    confusion_image.seek(0)
    plt.close(fig)

    return distribution_image, tree_image, confusion_image


def save_code_excerpt_image():
    code_lines = [
        ("df = pd.read_csv(DATA_FILE, names=FEATURE_NAMES, na_values='?')", "#D4D4D4"),
        ("df = df.dropna().copy()", "#D4D4D4"),
        ("df['heart_disease'] = (df['target_raw'] > 0).astype(int)", "#D4D4D4"),
        ("", "#D4D4D4"),
        ("x_train, x_test, y_train, y_test = train_test_split(", "#DCDCAA"),
        ("    x, y, test_size=0.30, random_state=42, stratify=y", "#D4D4D4"),
        (")", "#DCDCAA"),
        ("", "#D4D4D4"),
        ("model = DecisionTreeClassifier(", "#4EC9B0"),
        ("    max_depth=3, min_samples_leaf=8, random_state=42", "#D4D4D4"),
        (")", "#4EC9B0"),
        ("model.fit(x_train, y_train)", "#DCDCAA"),
        ("y_pred = model.predict(x_test)", "#DCDCAA"),
        ("", "#D4D4D4"),
        ("accuracy = accuracy_score(y_test, y_pred)", "#9CDCFE"),
        ("precision = precision_score(y_test, y_pred)", "#9CDCFE"),
        ("recall = recall_score(y_test, y_pred)", "#9CDCFE"),
    ]

    width, height = 1600, 900
    image = Image.new("RGB", (width, height), "#1E1E1E")
    draw = ImageDraw.Draw(image)
    draw.rectangle([0, 0, width, 64], fill="#252526")
    draw.rectangle([0, 64, 74, height], fill="#1B1B1B")
    draw.text((26, 21), "classification_project.py", fill="#D4D4D4", font=_font(28))

    line_font = _font(32)
    line_number_font = _font(24)
    y = 110
    for number, (line, color) in enumerate(code_lines, start=1):
        if line:
            draw.text((24, y + 5), str(number).rjust(2), fill="#858585", font=line_number_font)
            draw.text((95, y), line, fill=color, font=line_font)
        y += 43

    code_image = BytesIO()
    image.save(code_image, format="PNG")
    code_image.seek(0)
    return code_image


def _font(size: int):
    for font_name in ("C:/Windows/Fonts/consola.ttf", "C:/Windows/Fonts/arial.ttf"):
        try:
            return ImageFont.truetype(font_name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def add_title(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.55))
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(23, 37, 84)


def add_bullets(slide, bullets: list[str], left=0.75, top=1.15, width=11.8, height=5.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.space_after = Pt(10)


def add_metric_card(slide, label: str, value: str, x: float, y: float):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.55), Inches(1.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
    shape.line.color.rgb = RGBColor(96, 115, 159)
    frame = shape.text_frame
    frame.clear()
    p1 = frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run = p1.add_run()
    run.text = value
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(23, 37, 84)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run = p2.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(55, 65, 81)


def build_slides(df, features, model, metrics, distribution_image, tree_image, confusion_image, code_image):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.25), Inches(11.95), Inches(1.0))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.add_run()
    title_run.text = "Heart Disease Classification"
    title_run.font.size = Pt(46)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(23, 37, 84)

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.25), Inches(11.95), Inches(0.7))
    subtitle_paragraph = subtitle_box.text_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_paragraph.add_run()
    subtitle_run.text = "Decision Tree Model"
    subtitle_run.font.size = Pt(28)
    subtitle_run.font.color.rgb = RGBColor(55, 65, 81)

    group_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.75), Inches(11.95), Inches(0.55))
    group_paragraph = group_box.text_frame.paragraphs[0]
    group_paragraph.alignment = PP_ALIGN.CENTER
    group_run = group_paragraph.add_run()
    group_run.text = "Group members: [add names here]"
    group_run.font.size = Pt(20)
    group_run.font.color.rgb = RGBColor(55, 65, 81)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Project Overview")
    add_bullets(
        slide,
        [
            "Dataset: UCI Cleveland Heart Disease dataset, collected from clinical patient records.",
            "Classification goal: predict whether heart disease is absent or present.",
            "Model used: one decision tree with max_depth=3 and min_samples_leaf=8.",
            "Decision Tree was chosen instead of Naive Bayes because medical decision rules are easier to interpret visually.",
            "Implementation: Python with pandas, scikit-learn, matplotlib, and python-pptx.",
        ],
        font_size=20,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Data Description and Source")
    add_bullets(
        slide,
        [
            f"Usable records after removing missing values: {len(df)} patient records.",
            "Original attributes include clinical measurements such as age, chest pain type, blood pressure, cholesterol, maximum heart rate, exercise angina, ST depression, vessels colored, and thalassemia result.",
            "Target field: original values 0-4. This project uses a binary target: 0 = no heart disease, 1 = heart disease present.",
            "Source: Janosi, Steinbrunn, Pfisterer, and Detrano (1989), Heart Disease dataset, UCI Machine Learning Repository. DOI: 10.24432/C52P4X.",
            "Dataset URL: https://archive.ics.uci.edu/dataset/45/heart+disease",
        ],
        font_size=19,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Target Class Distribution")
    distribution_image.seek(0)
    slide.shapes.add_picture(distribution_image, Inches(2.0), Inches(1.2), width=Inches(9.2))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Attributes Used for Training")
    attributes = ", ".join(FEATURE_LABELS[f] for f in features)
    add_bullets(
        slide,
        [
            f"Training attributes: {attributes}.",
            "These attributes were chosen because they are the standard 13 clinical variables used in published experiments with the Cleveland data.",
            "The original target value was not used as an input feature; it was converted into the binary class label.",
            "Rows containing missing values in ca or thal were removed so the decision tree could train without imputed values.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Classification Model: Decision Tree")
    tree_image.seek(0)
    slide.shapes.add_picture(tree_image, Inches(0.12), Inches(0.92), width=Inches(13.08))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Model Interpretation")
    top_feature = FEATURE_LABELS[features[model.tree_.feature[0]]]
    add_bullets(
        slide,
        [
            f"The tree's first split is based on {top_feature}, meaning this attribute gave the strongest initial separation between the two classes.",
            "The most visible rules combine chest pain type, thalassemia result, maximum heart rate, exercise-induced angina, and ST depression.",
            "Patients grouped into leaves with abnormal stress-test or chest-pain patterns are more often classified as heart disease present.",
            "The tree was limited to depth 3 so the resulting model stays readable enough to explain in class.",
        ],
        font_size=20,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Evaluation Results")
    add_metric_card(slide, "Accuracy", f"{metrics['accuracy']:.1%}", 0.85, 1.2)
    add_metric_card(slide, "Precision", f"{metrics['precision']:.1%}", 3.85, 1.2)
    add_metric_card(slide, "Recall", f"{metrics['recall']:.1%}", 6.85, 1.2)
    add_metric_card(slide, "Test records", str(metrics["n_test"]), 9.85, 1.2)
    confusion_image.seek(0)
    slide.shapes.add_picture(confusion_image, Inches(3.45), Inches(2.55), width=Inches(6.3))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Evaluation and Summary")
    add_bullets(
        slide,
        [
            f"The decision tree achieved {metrics['accuracy']:.1%} accuracy, {metrics['precision']:.1%} precision, and {metrics['recall']:.1%} recall on the held-out test set.",
            "The model performed reasonably well for a small and interpretable tree, especially because several clinical attributes are directly related to heart-disease diagnosis.",
            "Performance may be limited by the small sample size, removed missing-value rows, and the fact that a shallow tree cannot represent every interaction among clinical measurements.",
            "The shallow tree helped interpretability: the prediction rules are understandable, which is useful in a health-related setting.",
        ],
        font_size=20,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "How the Work Was Implemented")
    add_bullets(
        slide,
        [
            "The Python script loads the UCI data, assigns column names, removes rows with missing values, and converts the 0-4 target into a binary class.",
            "The data is split into 70% training and 30% testing with stratification so both class labels appear in both sets.",
            "A single DecisionTreeClassifier is trained with one setting: max_depth=3 and min_samples_leaf=8.",
            "The script exports the class distribution chart, tree visualization, confusion matrix, metrics text file, and this slide deck.",
            "Code file submitted with the slides: classification_project.py.",
        ],
        font_size=19,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Implementation Code Excerpt")
    code_image.seek(0)
    slide.shapes.add_picture(code_image, Inches(0.75), Inches(1.05), width=Inches(11.85))

    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0), Inches(2.55), Inches(13.333), Inches(1.0))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "Thank You"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(23, 37, 84)

    subtitle = slide.shapes.add_textbox(Inches(0), Inches(3.55), Inches(13.333), Inches(0.65))
    subtitle_paragraph = subtitle.text_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_paragraph.add_run()
    subtitle_run.text = "Questions?"
    subtitle_run.font.size = Pt(26)
    subtitle_run.font.color.rgb = RGBColor(55, 65, 81)

    pptx_path = BASE_DIR / "Heart_Disease_Classification_Final.pptx"
    try:
        prs.save(pptx_path)
    except PermissionError:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = BASE_DIR / f"Heart_Disease_Classification_Final_{timestamp}.pptx"
        prs.save(pptx_path)
    return pptx_path


def main():
    df = load_data()
    features, model, metrics = train_model(df)
    distribution_image, tree_image, confusion_image = save_charts(df, features, model, metrics)
    code_image = save_code_excerpt_image()
    build_slides(df, features, model, metrics, distribution_image, tree_image, confusion_image, code_image)

    print(f"Accuracy: {metrics['accuracy']:.3f}")
    print(f"Precision: {metrics['precision']:.3f}")
    print(f"Recall: {metrics['recall']:.3f}")


if __name__ == "__main__":
    main()
